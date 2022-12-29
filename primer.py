import sqlite3
import sys
import docx
import pptx
import openpyxl


class Docx_file:

    def __init__(self):
        self.doc = docx.Document('Расписание.docx')

    def main(self, time, lenght, name, zal_name):
        self.doc.add_paragraph(
            f'Название фильма: {name} | Название зала: {zal_name} | Время сеанса {time} | Продолжительность сеанса {lenght}\n')
        self.doc.save('Расписание.docx')


class Xlsx_file:

    def __init__(self):
        self.fl = 'load schedule.xlsx'
        self.xls = openpyxl.load_workbook(self.fl)

    def main(self, time):
        ws = self.xls['Лист1']
        ws.append([time])
        self.xls.save(self.fl)


class Pptx_file:

    def __init__(self):
        self.ppt = pptx.Presentation()

    def main(self, time, lenght, name, zal_name):
        first_slide_layout = self.ppt.slide_layouts[0]
        slide = self.ppt.slides.add_slide(first_slide_layout)
        slide.shapes.title.text = f'Название фильма: {name}\nНазвание зала: {zal_name}\nВремя сеанса {time}\nПродолжительность сеанса {lenght}\n'
        self.ppt.save(f'Буклет {name}.pptx')


class Main_admin:

    def __init__(self):
        self.db = sqlite3.connect('kinoteatr.db')
        self.cur = self.db.cursor()

    def add_kinoteatr(self):
        name = input('Напишите название кинотеатра: ')
        adres = input('Напишите адрес кинотеатра: ')
        self.cur.execute('INSERT INTO kinoteatr(name, adres) VALUES(?, ?);', (name, adres))
        self.db.commit()
        print('Добавлено')

    def add_zal(self):
        name_kinoteatr = input('Напишите название кинотеатра: ')
        name = input('Напишите название зала: ')
        main.add_kresla(name, name_kinoteatr)

    def add_seans(self):
        zal_name = input('Напишите название зала: ')
        name = input('Напишите название сеанса: ')
        time = input('Напишите время начала сеанса: ')
        length = input('Напишите продолжительность сеанса: ')
        self.cur.execute('INSERT INTO seans(name, zal_name, time, length) VALUES(?, ?, ?, ?);',
                         (name, zal_name, time, length))
        self.db.commit()
        a = Docx_file()
        a.main(time, length, name, zal_name)
        a = Pptx_file()
        a.main(time, length, name, zal_name)
        print('Добавлено')

    def add_kresla(self, name, name_kinoteatr):
        row = input('Напишите количество рядов в зале: ')
        pole = input('Напишите количество кресел в каждом столбе: ')
        self.cur.execute('INSERT INTO zali_kinoteatr(name, name_kinoteatr, rows, poles) VALUES(?, ?, ?, ?);',
                         (name, name_kinoteatr, row, pole))
        self.db.commit()
        print('Добавлено')


class Main_user:

    def __init__(self):
        self.db = sqlite3.connect('kinoteatr.db')
        self.cur = self.db.cursor()

    def search_seans(self):
        c = self.cur.execute('SELECT name FROM seans').fetchall()
        for i in c:
            print(f'Фильм: {i[0]}')
        name = input('Введите название фильма: ')
        name = tuple([name])
        seans = self.cur.execute(f'SELECT * FROM seans WHERE name = ?', name)
        for i in seans:
            print(
                f'Название фильма: {i[1]} | Название зала: {i[2]} | Время сеанса {i[3]} | Продолжительность сеанса {i[4]}')
        time = input('Напишите время начала сеанса: ')
        name_zal = input('Напишите название зала: ')
        a = self.cur.execute(f'SELECT * FROM mesta_false WHERE name = ?', name).fetchall()
        b = self.cur.execute(f'SELECT * FROM zali_kinoteatr WHERE name = ?', tuple([name_zal])).fetchone()
        for i in range(b[3]):
            print(str(i + 1), end=' ')
            for j in range(b[4]):
                z = True
                for g in a:
                    if i == g[0] and j == g[1]:
                        print('X', end=' ')
                        z = False
                        break
                if z:
                    print('*', end=' ')
            print()
        row = input('Напишите ряд, на который вы хотите выбрать место: ')
        pole = input('Напишите номер места в ряду, который вы забронируете: ')
        self.cur.execute('INSERT INTO mesta_false(row, pole, name) VALUES(?, ?, ?)',
                         (int(row) - 1, int(pole) - 1, name[0]))
        self.db.commit()
        a = Xlsx_file()
        a.main(time)
        print('Забронировано')


login = input('Введите логин: ')
password = input('Введите пароль: ')
if login == 'admin' and password == 'admin':
    main = Main_admin()
    print('''Примеры команд:
    добавить кинотеатр
    добавить зал
    добавить сеанс\n-----------------------''')
    print('Напишите команду: ')
    for command in sys.stdin:
        if command == 'добавить кинотеатр\n':
            main.add_kinoteatr()
        elif command == 'добавить зал\n':
            main.add_zal()
        elif command == 'добавить сеанс\n':
            main.add_seans()

else:
    while True:
        main = Main_user()
        main.search_seans()
